from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# VI Impact brand
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_LINE = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def bar(slide, top=Inches(0.02)):
    s = slide.shapes.add_shape(1, 0, top, prs.slide_width, Inches(0.05))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def h1(slide, text, left=Inches(0.8), top=Inches(0.4), size=32):
    tb = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    return tb


def body(slide, text, left=Inches(0.8), top=Inches(1.6), w=Inches(5.5), size=18):
    tb = slide.shapes.add_textbox(left, top, w, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK if line.strip() else MUTED
        p.font.name = "Outfit"
        p.space_after = Pt(8 if line.strip() else 4)
    return tb


def bullets(slide, items, left=Inches(0.8), top=Inches(1.6), w=Inches(6), size=16):
    tb = slide.shapes.add_textbox(left, top, w, Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK if item.strip() else MUTED
        p.font.name = "Outfit"
        p.space_after = Pt(6 if item.strip() else 4)
    return tb


def card(slide, num, label, left, top, w=Inches(2.4), h=Inches(1.5)):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD_BG
    s.line.color.rgb = LIGHT_LINE
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED
    p2.font.name = "Outfit"
    p2.alignment = PP_ALIGN.CENTER


# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
# NOTE: python-pptx doesn't support SVG. Convert logo to PNG first:
#   pip install cairosvg && python3 -c "import cairosvg; cairosvg.svg2png(url='logo.svg', write_to='logo.png')"
# Or use Pillow to add PNG logo: s.shapes.add_picture('logo.png', left, top, w, h)
bar(s, top=Inches(2.6))
h1(s, "EDV Beratung Arnd Hekermans", left=Inches(2.5), top=Inches(2.85), size=34)
body(
    s,
    "Softwareentwicklung für Embedded Automotive Systeme\nVI Impact — Deterministic Traceability Engine",
    left=Inches(2.5),
    top=Inches(3.8),
    size=16,
)
# Contact line
tb = s.shapes.add_textbox(Inches(2.5), Inches(5.8), Inches(8), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "arnd@hekermans.de  ·  www.hekermans.de  ·  www.hekermans.de/vimpact.eu"
p.font.size = Pt(12)
p.font.color.rgb = MUTED
p.font.name = "Outfit"

# ===== SLIDE 2: PROFILE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Arnd Hekermans")
bullets(
    s,
    [
        "Dipl.-Inf. · Freelancer seit über 20 Jahren im Automotive-Bereich",
        "",
        "Fokus: Vernetzte Steuergeräte — Feldbusse, Protokolle, Applikationslogik",
        "",
        "Kernkompetenzen",
        "• Autosar Classic 3.x / 4.x — Diagnose (UDS, DCM), NVM, FEE, OS, ETH, FBL, DEM, MCAL",
        "• ISO 26262 Functional Safety bis ASIL-D · ASPICE v3",
        "• Signal-/Kontrollfluss-Analysen · Reverse Engineering · System-Dokumentation",
        "",
        "Kunden: Daimler, BMW, Bosch, Panasonic",
        "Interessen: Adaptive Autosar · BMS · ADAS · ISO 21434 Security",
    ],
)

# ===== SLIDE 3: VI OVERVIEW (Hero) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "VI Impact — The Traceability Engine")
body(
    s,
    "Connect the ALM tools you already run.\nSurface every gap. Pass the audit.",
    top=Inches(1.4),
    size=20,
)
body(
    s,
    "Ein Assessment-Layer über Ihrer Safety-Engineering-Toolchain.\nKein Ersatz — sondern die Schicht, die beweist,\ndass Ihre Toolchain produziert hat, was sie sollte.",
    top=Inches(2.4),
    size=15,
)
card(s, "9", "Adapters", Inches(0.8), Inches(5.0))
card(s, "4", "Tiers", Inches(3.5), Inches(5.0))
card(s, "6", "Gap Types", Inches(6.2), Inches(5.0))
card(s, "5+", "Norms", Inches(8.9), Inches(5.0))

# ===== SLIDE 4: PIPELINE (Stack) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Engine Architecture — 4 Stages")
steps = [
    (
        "01 · INGEST",
        "Adapters lesen aus Polarion, Jama,\ncodeBeamer, Git, MagicDraw, EA",
    ),
    ("02 · NORMALIZE", "Ein typisiertes Graph-Modell\nüber alle Tools hinweg"),
    (
        "03 · ENRICH",
        "Norm-Clauses (ISO 26262, 21434,\nASPICE, 13849) mit REQs verlinkt",
    ),
    ("04 · ANALYZE", "Deterministische Path-Walks,\nGap Detection, Coverage Reports"),
]
for i, (title, desc) in enumerate(steps):
    left = Inches(0.6 + i * 3.15)
    s2 = s.shapes.add_shape(1, left, Inches(1.8), Inches(2.9), Inches(1.8))
    s2.fill.solid()
    s2.fill.fore_color.rgb = CARD_BG
    s2.line.color.rgb = LIGHT_LINE
    s2.line.width = Pt(1)
    tf = s2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.font.name = "Outfit"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    # Arrow between boxes
    if i < 3:
        arrow = s.shapes.add_shape(
            6,
            left + Inches(2.95),
            Inches(2.45),
            Inches(0.25),
            Inches(0.3),  # right arrow
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# ===== SLIDE 5: GAP TYPES (Grid) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Six Gap Types — Classified by Severity")
gaps = [
    "Orphan\nRequirement",
    "Broken\nTrace",
    "Orphan\nTest",
    "Untraced\nCode",
    "Dead\nCode",
    "ASIL\nMismatch",
]
gap_desc = [
    "REQ ohne Link zu\nTest/Design/Norm",
    "Kette bricht\nan einer Stelle",
    "Test ohne\nRequirement",
    "Funktion ohne\nUpstream-REQ",
    "Keine Caller\nim Call-Graph",
    "Komplexität >\nASIL Class",
]
for i in range(6):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.8 + row * 2.5)
    s2 = s.shapes.add_shape(1, left, top, Inches(3.7), Inches(2.0))
    s2.fill.solid()
    s2.fill.fore_color.rgb = CARD_BG
    s2.line.color.rgb = LIGHT_LINE
    s2.line.width = Pt(1)
    tf = s2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = gaps[i]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = gap_desc[i]
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.font.name = "Outfit"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

# ===== SLIDE 6: WHY NOT AI (Split) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Why Not 'AI'?")
bullets(
    s,
    [
        "Compliance-Artefakte sind Evidence. Evidence muss reproduzierbar sein.",
        "",
        "Generativer Ansatz",
        "  LLM produziert finale Artefakte → LLM = Evidence",
        '  Auditor fragt: „Wie wurde entschieden?" → Antwort: Token-Verteilung',
        "  → Nicht auditierbar",
        "",
        "VI Ansatz",
        "  LLM schlägt vor → Engine validiert strukturell",
        "  → Engineer approved → Graph Walk = Evidence",
        "  → Gleicher Input = gleicher Output. Jeder Audit. Replayable.",
        "",
        "Der LLM ist nie die Source of Truth.",
    ],
    top=Inches(1.4),
    w=Inches(11.5),
    size=15,
)

# ===== SLIDE 7: DEPLOYMENT =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Deployment — Your Data, Your Network")
bullets(
    s,
    [
        "On-Premises: Container auf Kunden-Hardware · Single-Tenant",
        "Air-Gapped: Kein Netzwerk nach außen, kein Cloud-LLM, keine Telemetrie",
        "Hybrid RAG: Optionaler Chat-Endpoint mit Kunden-eigenem LLM",
        "",
        "Live-Deployments bei Tier-1 Zulieferern",
        "• Automotive ADAS  ·  Industrial Vehicles  ·  Process Automation",
        "• Concept Deployment ISO 13849: Jama + Git + AST/IR Layer",
        "• Mehrere tausend Items · Echte Norm-Corpora · Kunden-Hardware",
    ],
    top=Inches(1.6),
    w=Inches(11.5),
    size=16,
)

# ===== SLIDE 8: CONTACT (dark) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, PRIMARY)
h1(s, "Let's talk.", top=Inches(2.0), size=44)
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11), Inches(3))
tf = tb.text_frame
tf.word_wrap = True
lines = [
    "Arnd Hekermans",
    "arnd@hekermans.de  ·  +49 7153 926 0926",
    "www.hekermans.de  ·  www.hekermans.de/vimpact.eu",
    "",
    "EDV Beratung Arnd Hekermans",
    "Landskrona Weg 36  ·  73207 Plochingen",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18 if i == 0 else 14)
    p.font.bold = i == 0
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"
    p.space_after = Pt(6)
    p.alignment = PP_ALIGN.LEFT

# Override title color
for shape in s.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.font.size and p.font.size >= Pt(30):
                p.font.color.rgb = WHITE

output = "/Users/arndhekermans/Projekte/zed-fork/hekermans_vi_impact.pptx"
prs.save(output)
print(f"Saved: {output} ({len(prs.slides)} slides)")
