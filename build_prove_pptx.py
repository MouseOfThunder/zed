from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# VI Impact brand
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_LINE = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(s, c=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def bar(s, top=Inches(0.02)):
    sh = s.shapes.add_shape(1, 0, top, prs.slide_width, Inches(0.05))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()


def h1(s, text, left=Inches(0.8), top=Inches(0.4), size=32):
    tb = s.shapes.add_textbox(left, top, Inches(10), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"


def body(s, text, left=Inches(0.8), top=Inches(1.6), w=Inches(5.5), size=16):
    tb = s.shapes.add_textbox(left, top, w, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = "Outfit"
        p.font.color.rgb = DARK if line.strip() else MUTED
        p.space_after = Pt(8 if line.strip() else 4)


def bullets(s, items, left=Inches(0.8), top=Inches(1.6), w=Inches(11.5), size=15):
    tb = s.shapes.add_textbox(left, top, w, Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.name = "Outfit"
        p.font.color.rgb = DARK if item.strip() else MUTED
        p.space_after = Pt(6 if item.strip() else 4)


def card(s, num, label, left, top, w=Inches(2.4), h=Inches(1.5)):
    sh = s.shapes.add_shape(1, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD_BG
    sh.line.color.rgb = LIGHT_LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED
    p2.font.name = "Outfit"
    p2.alignment = PP_ALIGN.CENTER


def table_slide(s, rows, left=Inches(0.8), top=Inches(1.8), col_w=None):
    if col_w is None:
        col_w = [Inches(3.0), Inches(7.0)]
    tbl = s.shapes.add_table(
        len(rows), len(col_w), left, top, sum(col_w), Inches(0.45) * len(rows)
    ).table
    for ci in range(len(col_w)):
        tbl.columns[ci].width = col_w[ci]
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Outfit"
                p.font.color.rgb = PRIMARY if ri == 0 else DARK
                p.font.bold = ri == 0
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG


# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s, top=Inches(3.0))
h1(s, "PROVE", top=Inches(1.5), size=48)
body(
    s,
    "Process & Requirements Oversight & Verification Engine\n\nASPICE v4.0 Assessment Workbench\nby sensified Solutions GmbH",
    top=Inches(3.3),
    size=18,
)
body(s, "sensified.com  ·  Local LLM ·  Apple Silicon", top=Inches(5.5), size=13)

# ===== SLIDE 2: OVERVIEW =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "What is PROVE?")
body(
    s,
    "Bridges the gap between ALM tools and ASPICE assessments.\n\nIngests evidence → Builds knowledge graph\n→ Evaluates compliance → Generates reports.\n\nAll inference local. No cloud. Deterministic where it matters.",
    top=Inches(1.4),
    size=17,
)
card(s, "8", "Process\nGroups", Inches(0.8), Inches(5.0))
card(s, "18", "Evidence\nSources", Inches(3.5), Inches(5.0))
card(s, "3", "Assessment\nAngles", Inches(6.2), Inches(5.0))
card(s, "N-P-L-F", "ASPICE\nRating", Inches(8.9), Inches(5.0))

# ===== SLIDE 3: ARCHITECTURE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Architecture — 3 Services, 1 Command")
steps = [
    (
        "Frontend\n(port 5173)",
        "React + TypeScript\nVite dev server\nDashboard, Assessment, Solve",
    ),
    (
        "Backend API\n(port 8000)",
        "Python FastAPI\nGateway + MemoryGraph\nEvidence ingestion",
    ),
    ("LLM Server\n(port 8081)", "MLX / llama.cpp\nApple Silicon\nMetal acceleration"),
]
for i, (title, desc) in enumerate(steps):
    left = Inches(0.8 + i * 4.1)
    s2 = s.shapes.add_shape(1, left, Inches(2.0), Inches(3.7), Inches(2.2))
    s2.fill.solid()
    s2.fill.fore_color.rgb = CARD_BG
    s2.line.color.rgb = LIGHT_LINE
    s2.line.width = Pt(1)
    tf = s2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED
    p2.font.name = "Outfit"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    if i < 2:
        arrow = s.shapes.add_shape(
            6, left + Inches(3.75), Inches(2.8), Inches(0.25), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()
bullets(
    s,
    ["start.sh → all 3 services on localhost"],
    top=Inches(4.8),
    w=Inches(11.5),
    size=14,
)

# ===== SLIDE 4: EVIDENCE SOURCES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Evidence Ingestion — 18 Sources")
table_slide(
    s,
    [
        ["Tier", "Sources"],
        ["ALM / Requirements", "Polarion · DOORS NG · Jira · Confluence"],
        ["Source Code", "Git (MCU, NAD) · RTCU AST Analysis · ELF"],
        ["Test & Quality", "Robot Framework · Jira Bridge · CSV/JSON Imports"],
        [
            "Process Docs",
            "Confluence Management Evidence · Review Records · HLD/Architecture",
        ],
    ],
    left=Inches(0.8),
    top=Inches(1.6),
)

# ===== SLIDE 5: KNOWLEDGE GRAPH =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "MemoryGraph — Knowledge Graph")
bullets(
    s,
    [
        "Nodes: Requirements · Documents · Architecture Elements · Test Cases · Issues · Code Functions",
        "Edges: Implements · Traces · Calls · Depends On · Contains",
        "",
        "Capabilities:",
        "• FTS5 full-text search with diacritic removal",
        "• Automatic traceability gap detection",
        "• Coverage analysis across all linked evidence",
        "• Staleness detection — graph vs. database row count",
        "• Deterministic path walks → same input = same output",
    ],
    top=Inches(1.6),
)

# ===== SLIDE 6: ASPICE EVALUATION =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "ASPICE v4.0 Evaluation Engine")
table_slide(
    s,
    [
        ["Process Group", "Processes"],
        ["SYS — System", "SYS.2 · SYS.3"],
        ["SWE — Software", "SWE.1–SWE.6"],
        ["HWE — Hardware", "HWE.1–HWE.3"],
        ["MLE — ML Engineering", "MLE.1–MLE.3"],
        ["SEC — Security", "SEC.1–SEC.3"],
        ["MAN — Management", "MAN.3 · MAN.4 · MAN.5 · MAN.7"],
        ["SUP — Support", "SUP.1–SUP.8"],
        ["ACQ — Acquisition", "ACQ.1–ACQ.10"],
    ],
    left=Inches(0.8),
    top=Inches(1.6),
    col_w=[Inches(3.5), Inches(8.0)],
)
bullets(
    s,
    [
        "3 Assessment Angles:  PA 1.1 (Base Practices)  ·  PA 2.1 (Performance Mgmt)  ·  PA 2.2 (Work Products)",
        "Rating:   N (0-15%)  ·  P (16-50%)  ·  L (51-85%)  ·  F (86-100%)",
        "Capability Level 2:  PA 1.1 = F   AND   PA 2.1 ≥ L   AND   PA 2.2 ≥ L",
    ],
    top=Inches(6.0),
    w=Inches(11.5),
    size=13,
)

# ===== SLIDE 7: SOLVE MODE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "SOLVE Mode — Deterministic Engineering")
bullets(
    s,
    [
        "Complementary mode focused on gap resolution, not just detection.",
        "",
        "3 Domains:",
        "• Test Generation (SWE.4/SWE.5) — AST + libclang + Arrange/Oracle → NO LLM",
        "• Traceability — KG Linker + Name Matching → NO LLM",
        "• Requirements Quality — Quality Metrics + LLM Chat (local, optional)",
        "",
        "Key principle: Deterministic where it matters. LLM only where it helps.",
        "All inference on local MLX/llama.cpp — no cloud, no API keys.",
    ],
    top=Inches(1.6),
)

# ===== SLIDE 8: EXTENSION SYSTEM =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Extension System — Declarative JSON")
bullets(
    s,
    [
        "prove-extensions/ — fully configurable without code changes:",
        "",
        "• widgets/ — Dashboard widgets as JSON declarations",
        "• views/ — Custom page layouts",
        "• reports/ — Report templates for PDF, PPTX, HTML, Confluence",
        "• import-rules/ — File import mapping rules",
        "• layouts/ — Dashboard layout configurations",
        "• themes/ + logos/ — Customer branding (colors, fonts, SVG/PNG)",
        "",
        "Reports: Write-back to Confluence with automatic validation.",
        "Import: Rule-based ingestion from CSV, JSON, ELF, Robot Test outputs.",
    ],
    top=Inches(1.6),
)

# ===== SLIDE 9: TECH STACK =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
h1(s, "Technology Stack")
table_slide(
    s,
    [
        ["Layer", "Technology"],
        [
            "Frontend",
            "React 18 + TypeScript · Vite · Tailwind CSS · Ant Design · D3.js",
        ],
        ["Backend", "Python 3.11+ · FastAPI · SQLite (FTS5) · MemoryGraph"],
        ["LLM", "MLX (Apple Silicon) · llama.cpp · GGUF models — local only"],
        ["Extensions", "Declarative JSON configs · JSON Schema validation"],
        ["Output", "PDF (WeasyPrint) · PPTX (python-pptx) · HTML · Confluence REST"],
        ["Platform", "macOS · Apple Silicon optimized · Metal GPU acceleration"],
    ],
    left=Inches(0.8),
    top=Inches(1.6),
)

# ===== SLIDE 10: CONTACT =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, PRIMARY)
h1(s, "PROVE", top=Inches(2.0), size=44)
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11), Inches(3))
tf = tb.text_frame
tf.word_wrap = True
lines = [
    "Process & Requirements Oversight & Verification Engine",
    "",
    "sensified Solutions GmbH",
    "www.sensified.com",
    "",
    "ASPICE v4.0 Assessment Workbench",
    "Local LLM · Apple Silicon · Deterministic Engineering",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18 if i == 0 else 14)
    p.font.bold = i == 0
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"
    p.space_after = Pt(6)

for shape in s.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.font.size and p.font.size >= Pt(30):
                p.font.color.rgb = WHITE

output = "/Users/arndhekermans/Projekte/zed-fork/prove_overview.pptx"
prs.save(output)
print(f"Saved: {output} ({len(prs.slides)} slides)")
